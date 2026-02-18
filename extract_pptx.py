"""Extract text from PowerPoint presentation"""
from pptx import Presentation

pptx_file = "ADIVA-AI-Driven-Intelligent-Document-Validation-and-Analysis.pptx"

try:
    prs = Presentation(pptx_file)
    print(f"Total Slides: {len(prs.slides)}\n")
    print("=" * 80)
    
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n--- SLIDE {i} ---\n")
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(shape.text)
        
        print("\n" + "-" * 80)
    
    print(f"\n\nPresentation extracted successfully!")
    print(f"File: {pptx_file}")
    print(f"Total Slides: {len(prs.slides)}")

except Exception as e:
    print(f"Error: {e}")
